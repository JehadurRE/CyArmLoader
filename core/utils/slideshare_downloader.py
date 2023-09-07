import logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s [%(levelname)s]: %(message)s')
from io import BytesIO
import os
import img2pdf
import re
from os import walk
from os.path import join
import requests
from bs4 import BeautifulSoup
import pptx
import shutil
from PIL import Image 
from pytesseract import image_to_string
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib import colors
from reportlab.lib.units import inch, cm
from reportlab.pdfgen import canvas

CURRENT = os.path.dirname(__file__)


class SlideShareDownloader:
    """SlideShare Downloader Class"""

    def __init__(self, slideshare_url=None, download_format='pdf'):
        self.download_format = download_format
        self.slideshare_url = slideshare_url
        logging.info(f'Initialized with URL: {slideshare_url}, Format: {download_format}')

    def get_slide_info(self):
        try:
            html = requests.get(self.slideshare_url).content
            soup = BeautifulSoup(html, 'lxml')
            title = soup.find(class_='Heading_heading__LwpOS Heading_h1__J9yQZ Title_root__LXcGO').get_text().strip()
            source_tag = soup.find('source', attrs={'data-testid': 'slide-image-source'})
            srcset_value = source_tag['srcset']
            print(srcset_value)
            image_url = srcset_value
            print(len(image_url.split(',')))
            if(len(image_url.split(','))==3 ):
                final_img_url = image_url.split(',')[2].replace(
                ' ', '').replace('1024w', '')
            elif (len(image_url.split(','))==2 ):
                final_img_url = image_url.split(',')[1].split(' ')[1]
            else:
                final_img_url = image_url.split(',')[0].replace(
                ' ', '').replace('320w', '')
            total_slides = soup.find(class_='total-slides j-total-slides').get_text().strip()
            metadata = soup.find_all(class_='metadata-item')
            category = soup.find(class_='CategoryChips_root__6o2nr').get_text().strip()
            date = soup.find(class_='Text_root__Qdprv Text_secondary__SDKFB Text_medium__XbUIY').get_text().strip()
            views= soup.find_all(class_='Text_root__Qdprv Text_secondary__SDKFB Text_weight-strong__Cygpu Text_medium__XbUIY Likes_root__8tyVB')[3].get_text().strip()
            print(views)
            if len(metadata) >= 2:
                date, views = metadata[0].get_text(
                ).strip(), metadata[2].get_text().strip()

            return title, final_img_url, total_slides, category, date, views
        except Exception as e:
            logging.error(f"Failed to fetch slide information: {e}")
            print(f"Failed to fetch slide information: {e}")

    def get_file_name(self):
        # get url basename and replace non-alpha with '_'
        file_name = re.sub('[^0-9a-zA-Z]+', '_',
                           self.slideshare_url.split("/")[-1])
        if file_name.strip() == '':
            print(
                "Something wrong to get filename from URL, fallback to result.pdf or result.pptx")
            file_name = f"result.{self.download_format.lower()}"
        else:
            file_name += f".{self.download_format.lower()}"
        return file_name

    def download_images(self):
        try:
            html = requests.get(self.slideshare_url).content
            soup = BeautifulSoup(html, 'lxml')
            # soup.title.string
            title = '/tmp'
            images = soup.findAll('source', attrs={'data-testid': 'slide-image-source'})
            i = 0
            for image in images:
                image_url = image.get('srcset')
                print(image_url)
                length = len(image_url.split(','))
                final_img_url = image_url.split(',')[length-1].split(' ')[1]
                img = requests.get(final_img_url)
                if not os.path.exists(title):
                    os.makedirs(title)
                with open(f"{title}/{i}", 'wb') as f:
                    f.write(img.content)
                i += 1
            print(self.download_format)
            bfr, filename = self.convert(title)
            return bfr, filename
        except Exception as e:
            logging.error(f"Failed to download images: {e}")
            print(f"Failed to download images: {e}")

    def convert(self, img_dir_name):
        try:
            imgs = []
            for (dirpath, dirnames, filenames) in walk(join(CURRENT, img_dir_name)):
                imgs.extend(filenames)
                break
            imgs = ["%s/%s" % (img_dir_name, x) for x in imgs]

            def atoi(text):
                return int(text) if text.isdigit() else text

            def natural_keys(text):
                return [atoi(c) for c in re.split('(\d+)', text)]

            imgs.sort(key=natural_keys)

            f_bfr = BytesIO()
            filename = self.get_file_name()
            # if self.download_format == 'pdf':
            #     pdf_bytes = img2pdf.convert(imgs, dpi=300, x=None, y=None)
            #     with open(filename, "wb") as doc:
            #         doc.write(pdf_bytes)

            #     with open(filename, "rb") as fp:
            #         f_bfr.write(fp.read())
            #     f_bfr.write(pdf_bytes)
            if self.download_format == 'pdf':
                c = canvas.Canvas(filename, pagesize=letter)
                for img_path in imgs:
                    img = Image.open(img_path)
                    text = image_to_string(img)
                    c.drawString(100, 750, text)
                    c.showPage()
                c.save()
                
                with open(filename, "rb") as fp:
                    f_bfr.write(fp.read())
                
            else:
                p = pptx.Presentation()
                blank_slide_layout = p.slide_layouts[6]
                for im in imgs:
                    try:
                        slide = p.slides.add_slide(blank_slide_layout)
                        slide.shapes.add_picture(
                            im, 0, 0, p.slide_width, p.slide_height)
                    except Exception as e:
                        image = Image.open(im)
                        image = image.convert('RGB')
                        temp_image_path = f"temp_image.png"
                        image.save(temp_image_path, format='PNG')
                        slide.shapes.add_picture(
                        temp_image_path, 0, 0, p.slide_width, p.slide_height)
                        os.remove(temp_image_path)
                        logging.error(f"Conveted with exception image to pptx: 2048w")

                p.save(filename)
                with open(filename, "rb") as fp:
                    f_bfr.write(fp.read())

            f_bfr.seek(0)
            os.remove(filename)
            # shutil.rmtree(join(CURRENT, img_dir_name))
            directory_to_empty = join(CURRENT, img_dir_name)
            if os.path.exists(directory_to_empty) and os.path.isdir(directory_to_empty):
                # List all the files and subdirectories within the directory
                for item in os.listdir(directory_to_empty):
                    item_path = join(directory_to_empty, item)
                    
                    # Check if it's a file and delete it
                    if os.path.isfile(item_path):
                        os.remove(item_path)
            
            logging.info("Conversion successful.")
            return f_bfr, filename

        except Exception as e:
            logging.error(f"Conversion failed: {e}")
            print(f"Conversion failed: {e}")
            return None, None
